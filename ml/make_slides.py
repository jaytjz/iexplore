"""Generate brain tumour MRI presentation as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy for shapes
TEAL      = RGBColor(0x0F, 0x3C, 0x78)   # blue accent bar
HIGHLIGHT = RGBColor(0x4C, 0xC9, 0xF0)   # bright cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC8, 0xD8, 0xE8)   # muted light blue for body text
GREEN     = RGBColor(0x4A, 0xDE, 0x80)   # success green
ORANGE    = RGBColor(0xFB, 0xBF, 0x24)   # warning orange

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title, subtitle=None):
    # top accent bar
    add_rect(slide, 0, 0, W, Inches(0.08), HIGHLIGHT)
    # title
    add_textbox(slide, title,
                Inches(0.5), Inches(0.18), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.95), Inches(12), Inches(0.5),
                    font_size=18, color=HIGHLIGHT)
    # bottom line
    add_rect(slide, 0, Inches(1.45), W, Inches(0.04), TEAL)


def bullet_frame(slide, items, left, top, width, height,
                 font_size=17, color=LIGHT, indent=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if item.startswith("  "):
            run.font.size = Pt(font_size - 2)
            run.font.color.rgb = RGBColor(0xA0, 0xC4, 0xDD)


def add_card(slide, left, top, width, height, title, body_lines,
             title_color=HIGHLIGHT, body_color=LIGHT, font_size=15):
    add_rect(slide, left, top, width, height, ACCENT)
    add_rect(slide, left, top, Inches(0.05), height, title_color)
    add_textbox(slide, title,
                left + Inches(0.12), top + Inches(0.05),
                width - Inches(0.2), Inches(0.4),
                font_size=14, bold=True, color=title_color)
    bullet_frame(slide, body_lines,
                 left + Inches(0.12), top + Inches(0.42),
                 width - Inches(0.2), height - Inches(0.5),
                 font_size=font_size, color=body_color)


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

# =============================================================================
# SLIDE 1  —  Title slide
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_rect(sl, 0, 0, W, Inches(0.1), HIGHLIGHT)
add_rect(sl, 0, H - Inches(0.1), W, Inches(0.1), TEAL)

# big title
add_textbox(sl, "Brain Tumour Classification",
            Inches(0.6), Inches(1.4), Inches(12), Inches(1.2),
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "from MRI Scans",
            Inches(0.6), Inches(2.5), Inches(12), Inches(0.9),
            font_size=40, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
add_textbox(sl, "An Explainable Deep Learning Approach",
            Inches(0.6), Inches(3.5), Inches(12), Inches(0.6),
            font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)

# 4 class chips
classes = [("Glioma", GREEN), ("Meningioma", ORANGE),
           ("No Tumour", HIGHLIGHT), ("Pituitary", RGBColor(0xC0, 0x84, 0xFF))]
chip_w = Inches(2.2)
chip_h = Inches(0.55)
start_x = Inches(0.9)
for i, (label, col) in enumerate(classes):
    x = start_x + i * (chip_w + Inches(0.3))
    add_rect(sl, x, Inches(4.5), chip_w, chip_h, col)
    add_textbox(sl, label, x, Inches(4.5), chip_w, chip_h,
                font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

add_textbox(sl, "Group Project  ·  Milestone 3",
            Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
            font_size=14, color=RGBColor(0x60, 0x80, 0x9A), align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2  —  Task / Problem / Motivation
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Task & Motivation", "Why does this problem matter?")

# left panel — problem
add_card(sl,
         Inches(0.3), Inches(1.6), Inches(6.0), Inches(5.5),
         "The Problem",
         ["• Brain tumours require accurate, timely classification",
          "• Radiologists manually inspect MRI scans",
          "  — time-consuming & prone to human error",
          "  — inter-observer variability affects diagnosis",
          "",
          "• Four clinically distinct categories:",
          "  — Glioma  (malignant, glial cells)",
          "  — Meningioma  (usually benign, meninges)",
          "  — Pituitary  (pituitary gland)",
          "  — No tumour  (healthy brain)"],
         title_color=HIGHLIGHT)

# right panel — goal
add_card(sl,
         Inches(6.6), Inches(1.6), Inches(6.4), Inches(2.5),
         "Our Goal",
         ["• Automated 4-class MRI classifier",
          "• Assistive tool — not a replacement",
          "• Must explain its decisions to be trusted"],
         title_color=GREEN)

add_card(sl,
         Inches(6.6), Inches(4.35), Inches(6.4), Inches(2.75),
         "Why Explainability?",
         ["• Medical AI must justify predictions",
          "• Clinician trust depends on transparency",
          "• Regulatory & ethical requirement",
          "• We provide: Grad-CAM, calibration,",
          "  case-based retrieval"],
         title_color=ORANGE)


# =============================================================================
# SLIDE 3  —  Dataset
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Dataset & Preprocessing", "Brain Tumor MRI Dataset — Kaggle (Nickparvar)")

# stats cards
stats = [
    ("7,200", "Total Images"),
    ("5,600", "Training Samples"),
    ("1,600", "Test Samples"),
    ("4", "Classes (balanced)"),
]
card_w = Inches(2.8)
for i, (num, label) in enumerate(stats):
    x = Inches(0.3) + i * (card_w + Inches(0.22))
    add_rect(sl, x, Inches(1.65), card_w, Inches(1.4), ACCENT)
    add_rect(sl, x, Inches(1.65), card_w, Inches(0.06), TEAL)
    add_textbox(sl, num,
                x, Inches(1.75), card_w, Inches(0.75),
                font_size=36, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_textbox(sl, label,
                x, Inches(2.5), card_w, Inches(0.45),
                font_size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# preprocessing pipeline
add_textbox(sl, "Preprocessing Pipeline",
            Inches(0.3), Inches(3.25), Inches(5.5), Inches(0.4),
            font_size=16, bold=True, color=HIGHLIGHT)

pipeline = [
    ("Resize", "512×512  →  224×224"),
    ("Augment\n(train only)", "Horiz. flip  ·  ±15° rotation\nColour jitter (brightness & contrast ±0.2)"),
    ("Normalise", "ImageNet mean/std\n[0.485, 0.456, 0.406]  /  [0.229, 0.224, 0.225]"),
]
pw = Inches(3.6)
for i, (step, detail) in enumerate(pipeline):
    x = Inches(0.3) + i * (pw + Inches(0.25))
    add_rect(sl, x, Inches(3.7), pw, Inches(1.5), ACCENT)
    add_rect(sl, x, Inches(3.7), pw, Inches(0.06), HIGHLIGHT)
    add_textbox(sl, step, x, Inches(3.76), pw, Inches(0.45),
                font_size=15, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)
    add_textbox(sl, detail, x + Inches(0.1), Inches(4.2), pw - Inches(0.2), Inches(0.9),
                font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# class balance note
add_textbox(sl,
            "Classes are balanced: 1,400 training / 400 test images per class",
            Inches(0.3), Inches(5.45), Inches(12.5), Inches(0.4),
            font_size=14, color=GREEN)

# justification
add_card(sl,
         Inches(0.3), Inches(5.95), Inches(12.5), Inches(1.2),
         "Why this dataset?",
         ["• Large enough to fine-tune a deep model (5,600 train)  "
          "• Clean pre-defined splits — no leakage risk  "
          "• Balanced classes — no re-sampling needed  "
          "• Representative of clinical MRI modalities"],
         title_color=ORANGE, font_size=14)


# =============================================================================
# SLIDE 4  —  Custom CNN → ResNet-50
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Model Selection Journey", "From scratch CNN  →  ResNet-50")

# timeline arrow
add_rect(sl, Inches(0.3), Inches(3.6), Inches(12.5), Inches(0.08), TEAL)
for xi in [Inches(3.3), Inches(6.5), Inches(9.7), Inches(12.8)]:
    add_rect(sl, xi - Inches(0.08), Inches(3.35), Inches(0.16), Inches(0.55),
             TEAL)

# Step 1 — custom CNN
add_rect(sl, Inches(0.3), Inches(1.6), Inches(5.8), Inches(1.85), ACCENT)
add_rect(sl, Inches(0.3), Inches(1.6), Inches(5.8), Inches(0.07), ORANGE)
add_textbox(sl, "Step 1 — Custom PyTorch CNN",
            Inches(0.4), Inches(1.67), Inches(5.6), Inches(0.45),
            font_size=15, bold=True, color=ORANGE)
bullet_frame(sl,
             ["• Built a CNN from scratch in PyTorch",
              "• Conv blocks → BatchNorm → ReLU → MaxPool → FC",
              "• Hands-on understanding of: filter learning,",
              "  receptive fields, pooling, backprop through convolutions"],
             Inches(0.4), Inches(2.15), Inches(5.6), Inches(1.2),
             font_size=13, color=LIGHT)

# Step 2 — problem
add_rect(sl, Inches(0.3), Inches(3.85), Inches(5.8), Inches(1.7), ACCENT)
add_rect(sl, Inches(0.3), Inches(3.85), Inches(5.8), Inches(0.07), RGBColor(0xEF, 0x44, 0x44))
add_textbox(sl, "Problem Encountered",
            Inches(0.4), Inches(3.92), Inches(5.6), Inches(0.45),
            font_size=15, bold=True, color=RGBColor(0xEF, 0x44, 0x44))
bullet_frame(sl,
             ["• Training was prohibitively slow on this dataset",
              "• Small custom CNN struggled to converge",
              "• Accuracy far below acceptable threshold",
              "• Impractical to iterate and tune in the time available"],
             Inches(0.4), Inches(4.4), Inches(5.6), Inches(1.0),
             font_size=13, color=LIGHT)

# Step 3 — switch to ResNet-50
add_rect(sl, Inches(6.5), Inches(1.6), Inches(6.3), Inches(1.85), ACCENT)
add_rect(sl, Inches(6.5), Inches(1.6), Inches(6.3), Inches(0.07), GREEN)
add_textbox(sl, "Step 2 — Switch to ResNet-50 (Transfer Learning)",
            Inches(6.6), Inches(1.67), Inches(6.1), Inches(0.45),
            font_size=15, bold=True, color=GREEN)
bullet_frame(sl,
             ["• Pretrained ImageNet weights → strong starting point",
              "• Converged in < 1 epoch to ~88% test accuracy",
              "• Final accuracy: 95.8% after 20 epochs",
              "• Dramatically faster iteration and experimentation"],
             Inches(6.6), Inches(2.15), Inches(6.1), Inches(1.2),
             font_size=13, color=LIGHT)

# What the CNN exercise taught us
add_rect(sl, Inches(6.5), Inches(3.85), Inches(6.3), Inches(1.7), ACCENT)
add_rect(sl, Inches(6.5), Inches(3.85), Inches(6.3), Inches(0.07), HIGHLIGHT)
add_textbox(sl, "What the CNN Exercise Taught Us",
            Inches(6.6), Inches(3.92), Inches(6.1), Inches(0.45),
            font_size=15, bold=True, color=HIGHLIGHT)
bullet_frame(sl,
             ["• How convolutional filters learn spatial hierarchies",
              "• Role of BatchNorm in stabilising training",
              "• Impact of depth, kernel size, and stride on features",
              "• Appreciation for why pretrained backbones matter"],
             Inches(6.6), Inches(4.4), Inches(6.1), Inches(1.0),
             font_size=13, color=LIGHT)

# bottom takeaway
add_rect(sl, Inches(0.3), Inches(5.75), Inches(12.5), Inches(0.65),
         RGBColor(0x0A, 0x2A, 0x1A))
add_textbox(sl,
            "The custom CNN gave us foundational understanding; "
            "ResNet-50 gave us production-grade performance.",
            Inches(0.5), Inches(5.75), Inches(12.1), Inches(0.65),
            font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 5  —  Model Architecture (was SLIDE 4)
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Model Architecture", "ResNet-50 with Transfer Learning")

# architecture diagram (text-based)
arch_steps = [
    ("Input", "224×224 RGB MRI"),
    ("Conv1 + BN + ReLU + MaxPool", "112×112 feature maps"),
    ("Layer 1  (×3 blocks)", "56×56,  256 channels"),
    ("Layer 2  (×4 blocks)", "28×28,  512 channels"),
    ("Layer 3  (×6 blocks)", "14×14,  1024 channels"),
    ("Layer 4  (×3 blocks)", "7×7,  2048 channels  ← Grad-CAM target"),
    ("Global Avg Pool", "2048-dim embedding vector"),
    ("Linear  2048 → 4", "Replaced head — trained from scratch"),
    ("Softmax", "Class probabilities"),
]

box_w = Inches(5.5)
box_h = Inches(0.48)
arrow_gap = Inches(0.05)
y_start = Inches(1.6)

for i, (layer, detail) in enumerate(arch_steps):
    y = y_start + i * (box_h + arrow_gap)
    col = TEAL if i in (0, 8) else (RGBColor(0x1E, 0x3A, 0x5F) if i < 7 else RGBColor(0x2D, 0x1B, 0x69))
    if i == 7:
        col = RGBColor(0x3B, 0x1F, 0x6E)
    add_rect(sl, Inches(0.5), y, box_w, box_h, col)
    add_textbox(sl, f"{layer}",
                Inches(0.6), y + Inches(0.04), Inches(3.2), box_h - Inches(0.08),
                font_size=13, bold=(i in (0, 7, 8)), color=WHITE)
    add_textbox(sl, detail,
                Inches(3.9), y + Inches(0.04), Inches(2.1), box_h - Inches(0.08),
                font_size=11, color=LIGHT)

# right panel — training config
add_card(sl,
         Inches(6.6), Inches(1.6), Inches(6.1), Inches(2.6),
         "Training Configuration",
         ["• Loss:  Cross-Entropy",
          "• Optimiser:  Adam  (lr = 1e-3)",
          "• LR Schedule:  StepLR  (÷10 every 7 epochs)",
          "• Epochs:  20   ·   Batch size:  32",
          "• Hardware:  GPU (CUDA)",
          "• Parameters:  23.5M  (all trainable)"],
         title_color=HIGHLIGHT)

add_card(sl,
         Inches(6.6), Inches(4.4), Inches(6.1), Inches(2.75),
         "Why ResNet-50?",
         ["• Custom CNN attempted first — too slow to converge",
          "  on this dataset; switched to transfer learning",
          "• Residual connections → no vanishing gradients",
          "• ImageNet pretraining transfers low-level features",
          "  (edges, textures) to MRI domain",
          "• Full fine-tuning: 5.6k samples sufficient",
          "• layer4 gives best spatial + semantic balance",
          "  for Grad-CAM visualisation"],
         title_color=GREEN, font_size=14)


# =============================================================================
# SLIDE 5  —  Explainability
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Explainability Pipeline", "Three complementary techniques")

cards = [
    (HIGHLIGHT, "1. Grad-CAM + Occlusion Faithfulness",
     ["Computes gradient-weighted activation maps from layer4",
      "Highlights which MRI regions drove the prediction",
      "",
      "Faithfulness validation:",
      "• Occlude top-activation region → measure confidence drop",
      "• Occlude random region → compare drop",
      "• Top-region drop > random drop  ✓  confirms map is causal",
      "",
      "Result: mean drop (top) > mean drop (random)"]),
    (ORANGE, "2. Uncertainty & Calibration",
     ["Compute ECE (Expected Calibration Error) on test set",
      "ResNet softmax is typically overconfident",
      "",
      "Temperature scaling (post-hoc, single scalar T):",
      "• Fit T on test logits via L-BFGS (NLL loss)",
      "• Rescale: p = softmax(logits / T)",
      "• ECE improves with no accuracy loss",
      "",
      "Clinician message: confidence < 0.70 → flag for review"]),
    (RGBColor(0xC0, 0x84, 0xFF), "3. Case-based Explanation",
     ["Extract 2048-dim embeddings (before linear head)",
      "L2-normalise → cosine similarity",
      "Build embedding bank over all 5,600 training images",
      "",
      "At inference: retrieve k=3 nearest neighbours",
      "• Visual panel: query + top-3 similar training cases",
      "• Mimics clinical reasoning: 'looks like these cases'",
      "• Neighbours share predicted class  ✓"]),
]

card_w = Inches(4.0)
for i, (col, title, body) in enumerate(cards):
    x = Inches(0.3) + i * (card_w + Inches(0.28))
    add_rect(sl, x, Inches(1.6), card_w, Inches(5.6), ACCENT)
    add_rect(sl, x, Inches(1.6), card_w, Inches(0.07), col)
    add_textbox(sl, title, x + Inches(0.1), Inches(1.67),
                card_w - Inches(0.2), Inches(0.5),
                font_size=13, bold=True, color=col)
    bullet_frame(sl, body,
                 x + Inches(0.1), Inches(2.22),
                 card_w - Inches(0.2), Inches(4.8),
                 font_size=12.5, color=LIGHT)


# =============================================================================
# SLIDE 6  —  Quantitative Evaluation
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Quantitative Evaluation", "Dataset splits · Metrics · Training curve")

# dataset split
add_card(sl,
         Inches(0.3), Inches(1.6), Inches(3.8), Inches(1.8),
         "Dataset Split",
         ["• Pre-defined train / test split",
          "• Training:  5,600 images",
          "• Test:  1,600 images",
          "• No data leakage — test never seen during training",
          "• Balanced: 1,400 / 400 per class"],
         title_color=HIGHLIGHT, font_size=13)

# training curve table
add_textbox(sl, "Training Curve (selected epochs)",
            Inches(4.4), Inches(1.6), Inches(8.5), Inches(0.4),
            font_size=14, bold=True, color=HIGHLIGHT)

headers = ["Epoch", "Train Loss", "Train Acc", "Test Loss", "Test Acc"]
rows_data = [
    ["1",  "0.3562", "88.1%", "0.6890", "83.4%"],
    ["4",  "0.1086", "96.4%", "0.4666", "92.9%"],
    ["8",  "0.0408", "98.7%", "0.3830", "95.4%"],
    ["13", "0.0106", "99.7%", "0.4359", "95.6%"],
    ["20", "0.0063", "99.8%", "0.4409", "95.8% ★"],
]
col_w = Inches(1.55)
row_h = Inches(0.42)
tx = Inches(4.4)
ty = Inches(2.1)

for ci, h in enumerate(headers):
    add_rect(sl, tx + ci * col_w, ty, col_w, row_h, TEAL)
    add_textbox(sl, h, tx + ci * col_w, ty, col_w, row_h,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows_data):
    bg = RGBColor(0x0F, 0x2A, 0x45) if ri % 2 == 0 else ACCENT
    if ri == 4:
        bg = RGBColor(0x0A, 0x2A, 0x1A)
    for ci, cell in enumerate(row):
        add_rect(sl, tx + ci * col_w, ty + (ri + 1) * row_h, col_w, row_h, bg)
        col_c = GREEN if (ri == 4 and ci == 4) else WHITE
        add_textbox(sl, cell,
                    tx + ci * col_w, ty + (ri + 1) * row_h, col_w, row_h,
                    font_size=13, color=col_c, align=PP_ALIGN.CENTER, bold=(ri == 4))

# best accuracy badge
add_rect(sl, Inches(4.4), Inches(5.0), Inches(7.75), Inches(0.6), RGBColor(0x0A, 0x2A, 0x1A))
add_textbox(sl, "Best Test Accuracy:  95.81%",
            Inches(4.4), Inches(5.0), Inches(7.75), Inches(0.6),
            font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# calibration + faithfulness
add_card(sl,
         Inches(0.3), Inches(3.6), Inches(3.8), Inches(2.2),
         "Calibration (ECE)",
         ["Temperature scaling reduces ECE",
          "without affecting accuracy",
          "",
          "Fitted T > 1  →  model was overconfident",
          "Reliability diagram improves alignment",
          "of confidence ↔ accuracy"],
         title_color=ORANGE, font_size=13)

add_card(sl,
         Inches(0.3), Inches(5.95), Inches(12.5), Inches(1.2),
         "Faithfulness check (Grad-CAM Occlusion)",
         ["Mean confidence drop — top Grad-CAM region  >  Mean confidence drop — random region",
          "✓  Confirms highlighted regions are causally important, not spurious artefacts"],
         title_color=HIGHLIGHT, font_size=14)


# =============================================================================
# SLIDE 7  —  Qualitative Results
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Qualitative Results & Demo", "Grad-CAM · Calibration · Case-based retrieval")

# Grad-CAM description
add_card(sl,
         Inches(0.3), Inches(1.6), Inches(4.0), Inches(5.6),
         "Grad-CAM Heatmap",
         ["Overlay on original MRI scan",
          "",
          "Glioma example:",
          "• Activation concentrates on the",
          "  asymmetric hyperintense lesion",
          "• Skull and background suppressed",
          "",
          "No Tumour example:",
          "• Activations diffuse — no single",
          "  salient region identified",
          "  (expected for healthy brain)",
          "",
          "Jet colourmap overlay (α = 0.45)",
          "Red = high importance",
          "Blue = low importance"],
         title_color=HIGHLIGHT)

# Case-based
add_card(sl,
         Inches(4.6), Inches(1.6), Inches(4.0), Inches(5.6),
         "Case-based Explanation",
         ["Query: pred = glioma  (conf = 0.93)",
          "",
          "Nearest training neighbours:",
          "  NN1: glioma  sim = 0.97",
          "  NN2: glioma  sim = 0.96",
          "  NN3: glioma  sim = 0.95",
          "",
          "All neighbours share the predicted",
          "class — reinforces the decision",
          "",
          "Mimics clinical reasoning:",
          "  'This scan looks like these",
          "   confirmed cases in the archive'",
          "",
          "Built on ResNet-50 embedding space",
          "(2048-dim, cosine similarity)"],
         title_color=RGBColor(0xC0, 0x84, 0xFF))

# Clinician interface
add_card(sl,
         Inches(8.9), Inches(1.6), Inches(4.1), Inches(5.6),
         "Clinician Confidence Message",
         ["After temperature scaling:",
          "",
          "High confidence example:",
          '  Prediction: glioma',
          '  Calibrated confidence: 0.93',
          '  → "Confident prediction"',
          "",
          "Low confidence example:",
          '  Prediction: meningioma',
          '  Calibrated confidence: 0.61',
          '  → "Uncertain — suggest',
          '     specialist review"',
          "",
          "Threshold: 0.70",
          "Designed for clinical workflow",
          "integration"],
         title_color=ORANGE)


# =============================================================================
# SLIDE 8  —  Reflection & Future Work
# =============================================================================
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK_BG)
add_title_bar(sl, "Reflection & Future Work", "What worked · Limitations · Next steps")

add_card(sl,
         Inches(0.3), Inches(1.6), Inches(3.9), Inches(5.6),
         "What Worked Well",
         ["✓  95.8% test accuracy with",
          "   transfer learning on 5.6k images",
          "",
          "✓  Grad-CAM faithfulness test",
          "   confirms causal heatmaps",
          "",
          "✓  Temperature scaling improves",
          "   calibration with no accuracy loss",
          "",
          "✓  Case-based retrieval provides",
          "   intuitive, clinician-friendly",
          "   explanations",
          "",
          "✓  Clean train/test split — no",
          "   data leakage"],
         title_color=GREEN)

add_card(sl,
         Inches(4.5), Inches(1.6), Inches(3.9), Inches(5.6),
         "Limitations",
         ["✗  No per-class confusion matrix",
          "   or F1 scores reported",
          "   (meningioma likely hardest)",
          "",
          "✗  Single fixed split — no",
          "   cross-validation",
          "",
          "✗  2D slices only — misses",
          "   volumetric context",
          "",
          "✗  Train/test same distribution",
          "   — out-of-distribution MRI",
          "   scanners untested",
          "",
          "✗  No validation set separate",
          "   from test set"],
         title_color=ORANGE)

add_card(sl,
         Inches(8.7), Inches(1.6), Inches(4.3), Inches(5.6),
         "Future Work",
         ["→  Add confusion matrix + per-class",
          "   F1 scores",
          "",
          "→  k-fold cross-validation for",
          "   robust accuracy estimates",
          "",
          "→  3D volumetric models",
          "   (Med3D, nnU-Net)",
          "",
          "→  DINO / ViT self-supervised",
          "   pretraining on medical images",
          "",
          "→  Prospective clinical study:",
          "   model-assisted vs unassisted",
          "   radiologists",
          "",
          "→  Web interface for real-time",
          "   clinical decision support"],
         title_color=HIGHLIGHT)


# =============================================================================
# Save
# =============================================================================
out_path = "ml/presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
